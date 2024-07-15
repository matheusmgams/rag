import os
from pptx import Presentation
from langchain.docstore.document import Document
from langchain_community.document_loaders import WebBaseLoader, PyMuPDFLoader

class DocumentLoader:
    def __init__(self, archives, codes, websites):
        # Define defaults para lidar com ausência de variáveis
        if not archives:
            print("O diretório de arquivos (archives) não foi definido. Documentos não serão carregados.")
            archives = None
        if not codes:
            print("A extensão dos arquivos de código (codes) não foi definida. Arquivos de código não serão carregados.")
            codes = None
        if not websites:
            print("A lista de websites (websites) não foi definida. Websites não serão carregados.")
            websites = []
        
        self.archives = archives
        self.codes = codes
        self.websites = websites

    def load_codes(self):
        if not self.archives or not self.codes:
            return []

        all_code = []
        for root, dirs, files in os.walk(self.archives):
            code_files = [os.path.join(root, file) for file in files if file.endswith(self.codes)]
            for code_file in code_files:
                try:
                    with open(code_file, 'r', encoding='utf-8') as f:
                        code = f.read()
                        all_code.append(Document(page_content=code, metadata={"title": code_file}))
                        print(f"Load {code_file}")
                except Exception as e:
                    print(f"Erro ao carregar {code_file}: {e}")
        return all_code

    def load_pdfs(self):
        if not self.archives:
            return []

        if not os.path.isdir(self.archives):
            raise FileNotFoundError(f"O diretório de arquivos '{self.archives}' não foi encontrado.")
        
        pdf_files = [os.path.join(self.archives, code_file) for code_file in os.listdir(self.archives) if code_file.endswith('.pdf')]
        all_pdfs = []
        for pdf_file in pdf_files:
            try:
                loader = PyMuPDFLoader(pdf_file)
                docs = loader.load()
                for doc in docs:
                    doc.metadata["title"] = pdf_file
                all_pdfs.extend(docs)
                print(f"Load {pdf_file}")
            except Exception as e:
                print(f"Erro ao carregar {pdf_file}: {e}")
        return all_pdfs

    def load_pptx(self):
        if not self.archives:
            return []

        if not os.path.isdir(self.archives):
            raise FileNotFoundError(f"O diretório de arquivos '{self.archives}' não foi encontrado.")
        
        pptx_files = [os.path.join(self.archives, code_file) for code_file in os.listdir(self.archives) if code_file.endswith('.pptx')]
        all_pptx = []
        for pptx_file in pptx_files:
            try:
                doc_content = ""
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            doc_content += shape.text + "\n"
                all_pptx.append(Document(page_content=doc_content, metadata={"title": pptx_file}))
                print(f"Load {pptx_file}")
            except Exception as e:
                print(f"Erro ao carregar {pptx_file}: {e}")
        return all_pptx

    def load_websites(self):
        if not self.websites:
            return []

        all_sites = []
        for website in self.websites:
            try:
                loader = WebBaseLoader(website)
                docs = loader.load()
                all_sites.extend(docs)
                print(f"Load {website}")
            except Exception as e:
                print(f"Erro ao carregar {website}: {e}")
        return all_sites

    def load_all_documents(self):
        all_docs = []
        if self.archives is not None and self.codes is not None:
            all_docs.extend(self.load_codes())
        if self.archives is not None:
            all_docs.extend(self.load_pdfs())
        if self.archives is not None:
            all_docs.extend(self.load_pptx())
        if self.websites is not None:
            all_docs.extend(self.load_websites())
        return all_docs
