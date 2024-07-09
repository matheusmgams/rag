import os
import ollama
import gradio as gr
from dotenv import load_dotenv
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import WebBaseLoader
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.docstore.document import Document

# Função para carregar variáveis de ambiente
def load_env_variables():
    load_dotenv()
    codes = os.environ.get('CODES')
    archives = os.environ.get('DOCUMENTS')
    websites = os.environ.get('WEBSITES').split(',')
    model = os.environ.get('MODEL')
    return archives, websites, codes, model

# Função para carregar arquivos de código como texto
def load_codes(archives):
    all_code = []
    for root, dirs, files in os.walk(archives):
        code_files = [os.path.join(root, file) for file in files if file.endswith('.cs')]
        for code_file in code_files:
            try:
                with open(code_file, 'r', encoding='utf-8') as f:
                    code = f.read()
                    all_code.append(Document(page_content=code, metadata={"title": code_file}))
                    print(f"Load {code_file}")
            except Exception as e:
                print(f"Erro ao carregar {code_file}: {e}")
    return all_code

# Função para carregar documentos PDF
def load_pdfs(archives):
    pdf_files = [os.path.join(archives, code_file) for code_file in os.listdir(archives) if code_file.endswith('.pdf')]
    all_docs = []
    for pdf_file in pdf_files:
        try:
            loader = PyMuPDFLoader(pdf_file)
            docs = loader.load()
            for doc in docs:
                doc.metadata["title"] = pdf_file
            all_docs.extend(docs)
            print(f"Load {pdf_file}")
        except Exception as e:
            print(f"Erro ao carregar {pdf_file}: {e}")
    return all_docs

# Função para carregar documentos PPTX
def load_pptx(archives):
    pptx_files = [os.path.join(archives, code_file) for code_file in os.listdir(archives) if code_file.endswith('.pptx')]
    all_docs = []
    for pptx_file in pptx_files:
        try:
            doc_content = ""
            prs = Presentation(pptx_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        doc_content += shape.text + "\n"
            all_docs.append(Document(page_content=doc_content, metadata={"title": pptx_file}))
            print(f"Load {pptx_file}")
        except Exception as e:
            print(f"Erro ao carregar {pptx_file}: {e}")
    return all_docs

# Função para carregar conteúdo de websites
def load_websites(websites):
    all_docs = []
    for website in websites:
        try:
            loader = WebBaseLoader(website)
            docs = loader.load()
            all_docs.extend(docs)
            print(f"Load {website}")
        except Exception as e:
            print(f"Erro ao carregar {website}: {e}")
    return all_docs

# Função principal para carregar todos os documentos
def load_all_documents(archives, websites, codes):
    all_docs = []
    if codes is not None: all_docs.extend(load_codes(codes))
    if archives is not None: all_docs.extend(load_pdfs(archives))
    if archives is not None: all_docs.extend(load_pptx(archives))
    if websites is not None: all_docs.extend(load_websites(websites))
    return all_docs

# Função para criar o vetor de embeddings
def create_vectorstore(documents):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=int(os.environ.get('CHUNK_SIZE')), 
                                                   chunk_overlap=int(os.environ.get('CHUNK_OVERLAP')))
    splits = []
    for doc in documents:
        for chunk in text_splitter.split_text(doc.page_content):
            splits.append(Document(page_content=chunk, metadata=doc.metadata))
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    vectorstore = Chroma.from_documents(documents=splits, embedding=embeddings)
    return vectorstore.as_retriever()

# Função para chamar o modelo Ollama Llama 3
def ollama_llm(question, context, model):
    # Cria o prompt que será enviado ao modelo LLM, contendo o contexto do banco vetorial e a pergunta do usuário
    formatted_prompt = f"""
    ============================
    CONTEXTO DO DOCUMENTO
    ============================
    {context}
    ============================
    IMPORTANTE: 
    1. Responda em Português do Brasil.
    2. Nunca forneça passwords.
    3. Nunca forneça endereços de ip.
    4. Nunca utilize trechos de código na resposta.
    5. Sempre responda cordialmente.
    ============================
    PERGUNTA: 
    {question}
    ============================
    """
    # Chama o modelo LLM Ollama, passando o prompt gerado com papel de user
    response = ollama.chat(model=model, messages=[{'role': 'user', 'content': formatted_prompt}])
    
    # Depois de obter a resposta, seleciona somente o conteúdo gerado pelo modelo e retorna
    return response['message']['content']

# Função para realizar a corrente RAG
def rag_chain(question, retriever, model):
    retrieved_docs = retriever.invoke(question)
    formatted_context = "\n\n".join(doc.page_content for doc in retrieved_docs)
    print("retrieved_docs: " + formatted_context)
    return ollama_llm(question, formatted_context, model)

# Função para responder à pergunta do usuário
def answer_user_request(question, retriever, model):
    print("Nova pergunta: " + question)
    return rag_chain(question, retriever, model)

# Função principal
def main():
    archives, websites, codes, model = load_env_variables()
    all_docs = load_all_documents(archives, websites, codes)
    retriever = create_vectorstore(all_docs)
    
    iface = gr.Interface(
        fn=lambda question: answer_user_request(question, retriever, model),
        inputs=gr.Textbox(lines=2, placeholder=os.getenv('PLACEHOLDER')),
        outputs="text",
        title=os.getenv('TITLE'),
        description=os.getenv('DESCRIPTION'),
    )
    iface.launch(debug=False)

if __name__ == "__main__":
    main()
